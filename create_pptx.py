from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util

# Brand colours
ORANGE   = RGBColor(0xFF, 0x6B, 0x1A)   # CampusTrade orange
DARK     = RGBColor(0x1A, 0x1A, 0x2E)   # Deep navy
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFF)
GRAY     = RGBColor(0x64, 0x64, 0x80)
GREEN    = RGBColor(0x10, 0xB9, 0x81)
RED      = RGBColor(0xEF, 0x44, 0x44)
BLUE     = RGBColor(0x1A, 0x56, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ── helpers ────────────────────────────────────────────────────────────────

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height, text, font_size=18,
        bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        bg_color=None, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    return txBox

def rect(slide, left, top, width, height, color, radius=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def bullet_box(slide, left, top, width, height, title, bullets,
               title_color=ORANGE, bullet_color=DARK, title_size=20,
               bullet_size=15, bg_color=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    # title paragraph
    tp = tf.paragraphs[0]
    tp.alignment = PP_ALIGN.LEFT
    tr = tp.add_run()
    tr.text = title
    tr.font.bold = True
    tr.font.size = Pt(title_size)
    tr.font.color.rgb = title_color
    # bullet paragraphs
    for b in bullets:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = "  • " + b
        r.font.size = Pt(bullet_size)
        r.font.color.rgb = bullet_color
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    return txBox

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
rect(s, 0, 0, 13.33, 0.12, ORANGE)
rect(s, 0, 7.38, 13.33, 0.12, ORANGE)

box(s, 0.8, 0.3, 11.7, 0.7,
    "SEN3244 — SOFTWARE ARCHITECTURE  |  ICT UNIVERSITY  |  SPRING 2026",
    font_size=11, color=ORANGE, align=PP_ALIGN.CENTER)

box(s, 1.2, 1.1, 10.9, 1.5,
    "CampusTrade", font_size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

box(s, 1.2, 2.55, 10.9, 0.6,
    "Smart Campus Marketplace", font_size=26, color=ORANGE, align=PP_ALIGN.CENTER)

box(s, 1.2, 3.2, 10.9, 0.5,
    "A production-grade event-driven microservices platform for ICT University students",
    font_size=16, color=RGBColor(0xCC, 0xCC, 0xEE), align=PP_ALIGN.CENTER)

# Stat boxes
for i, (num, lbl) in enumerate([("9", "Microservices"), ("232", "Tests Passing"),
                                  ("95%+", "Code Coverage"), ("13", "CI/CD Stages")]):
    x = 1.2 + i * 2.8
    rect(s, x, 4.0, 2.4, 1.3, RGBColor(0x2A, 0x2A, 0x4E))
    box(s, x, 4.05, 2.4, 0.65, num, font_size=30, bold=True,
        color=ORANGE, align=PP_ALIGN.CENTER)
    box(s, x, 4.65, 2.4, 0.45, lbl, font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

box(s, 1.2, 5.6, 10.9, 0.45,
    "Praises Ncha (Scrum Master)  ·  Kongyu Jesse Ntani (Product Owner)",
    font_size=15, color=RGBColor(0xCC, 0xCC, 0xEE), align=PP_ALIGN.CENTER)

box(s, 1.2, 6.1, 10.9, 0.35,
    "http://209.38.199.108:4000  |  github.com/Nkinyampraises/smart-campus-market",
    font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

box(s, 1.2, 6.5, 10.9, 0.35,
    "Instructor: Engr. TEKOH PALMA",
    font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Team
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_BG)
rect(s, 0, 0, 13.33, 1.0, DARK)
box(s, 0.5, 0.2, 12.3, 0.65, "Meet the Team", font_size=30, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)

for i, (name, role, badge, items, col) in enumerate([
    ("Praises Ncha", "Scrum Master", "SCRUM MASTER",
     ["Design & Frontend (React, Tailwind)", "Unit & Integration Testing (Jest)",
      "Jenkins CI/CD Pipeline (13 stages)", "VPS Deployment (DigitalOcean)",
      "Monitoring: Prometheus + Grafana", "SonarQube Quality Gates"], ORANGE),
    ("Kongyu Jesse Ntani", "Product Owner  |  ICTU20234195", "PRODUCT OWNER",
     ["Backend Microservices (9 services)", "Event-Driven Architecture (Redis)",
      "API Design & Database Schema", "Kubernetes Manifests & HPA",
      "Architecture Documentation", "PowerPoint & Presentation"], BLUE),
]):
    x = 0.5 + i * 6.4
    rect(s, x, 1.15, 6.0, 5.9, WHITE)
    rect(s, x, 1.15, 6.0, 0.5, col)
    box(s, x + 0.1, 1.15, 5.8, 0.45, badge, font_size=13, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    box(s, x + 0.1, 1.7, 5.8, 0.45, name, font_size=20, bold=True,
        color=DARK, align=PP_ALIGN.CENTER)
    box(s, x + 0.1, 2.1, 5.8, 0.4, role, font_size=13, color=GRAY,
        align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        box(s, x + 0.2, 2.6 + j * 0.7, 5.6, 0.55, "• " + item,
            font_size=14, color=DARK)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Problem Statement
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
rect(s, 0, 0, 13.33, 1.0, RGBColor(0x2A, 0x0A, 0x00))
box(s, 0.5, 0.18, 12.3, 0.65, "The Problem", font_size=32, bold=True,
    color=ORANGE, align=PP_ALIGN.CENTER)

box(s, 0.5, 1.1, 12.3, 0.45,
    "ICT University students trade via WhatsApp — no structure, no trust, no safety.",
    font_size=17, color=RGBColor(0xCC, 0xCC, 0xEE), align=PP_ALIGN.CENTER)

problems = [
    ("🚫", "No Trusted Platform",    "WhatsApp messages with zero verification or fraud protection"),
    ("💰", "Price Exploitation",      "New students overpay — no reference for fair campus pricing"),
    ("💬", "No Real-Time Negotiation","Must switch to external apps, losing listing context"),
    ("🔍", "No Searchable Inventory", "Listings vanish in chat history — impossible to discover"),
    ("⚠️",  "No Admin Oversight",      "Fraudulent sellers and spam listings go unchecked"),
]
for i, (icon, title, desc) in enumerate(problems):
    x = 0.4 + (i % 3) * 4.2
    y = 1.8 if i < 3 else 4.2
    if i == 3: x = 2.25
    if i == 4: x = 6.45
    rect(s, x, y, 3.8, 1.9, RGBColor(0x2A, 0x2A, 0x4E))
    box(s, x + 0.1, y + 0.1, 0.7, 0.7, icon, font_size=22, color=WHITE)
    box(s, x + 0.8, y + 0.1, 2.8, 0.45, title, font_size=14, bold=True,
        color=ORANGE)
    box(s, x + 0.1, y + 0.65, 3.5, 1.0, desc, font_size=12,
        color=RGBColor(0xCC, 0xCC, 0xEE))

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Our Solution
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_BG)
rect(s, 0, 0, 13.33, 1.0, DARK)
box(s, 0.5, 0.18, 12.3, 0.65, "Our Solution — CampusTrade", font_size=30,
    bold=True, color=WHITE, align=PP_ALIGN.CENTER)

solutions = [
    (GREEN,  "Structured Marketplace",    "Post listings with photos, categories, campus zones, and descriptions"),
    (ORANGE, "AI Fraud Detection",        "Auto-flags suspicious prices; red warning shown to all buyers"),
    (BLUE,   "Real-Time Chat",            "Socket.IO messaging — negotiate directly on the platform"),
    (GREEN,  "Push Notifications",        "Browser push + email for every offer, message, and event"),
    (ORANGE, "Admin Moderation",          "Full dashboard: suspend users, remove listings, resolve reports"),
    (BLUE,   "Full-Text Search",          "Find anything by keyword, category, price range, campus zone"),
]
for i, (col, title, desc) in enumerate(solutions):
    x = 0.4 + (i % 2) * 6.3
    y = 1.2 + (i // 2) * 1.85
    rect(s, x, y, 6.0, 1.6, WHITE)
    rect(s, x, y, 0.18, 1.6, col)
    box(s, x + 0.35, y + 0.15, 5.5, 0.45, title, font_size=16, bold=True, color=DARK)
    box(s, x + 0.35, y + 0.6,  5.5, 0.85, desc,  font_size=13, color=GRAY)

box(s, 0.5, 7.0, 12.3, 0.35,
    "Live at: http://209.38.199.108:4000",
    font_size=13, color=ORANGE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Architecture Overview
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
rect(s, 0, 0, 13.33, 1.0, RGBColor(0x0A, 0x0A, 0x1E))
box(s, 0.5, 0.18, 12.3, 0.65, "System Architecture", font_size=32, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)

arch = """
   [Student Browser]          [Admin Browser]
          │                         │
   ┌──────▼─────────────────────────▼──────┐
   │      NGINX  ·  Port 4000  (React SPA)  │
   └──────────────────┬────────────────────┘
                      │  HTTP/REST
   ┌──────────────────▼────────────────────┐
   │         API GATEWAY  ·  Port 8080      │
   │   JWT Validation  ·  Rate Limiting     │
   └──┬────┬────┬────┬────┬────┬────┬─────┘
      │    │    │    │    │    │    │
   Auth  User  List Chat Admin  AI  Search Notify
   3001  3002  3003  3004  3005 3006  3007   3008
      │    │    │    │    │    │    │    │
   ┌──▼────▼────▼────▼────▼────▼────▼────▼──┐
   │  PostgreSQL :5432  │  Redis :6379        │
   │  (Persistent data) │  (Cache + Pub/Sub)  │
   └───────────────────────────────────────┘
"""
box(s, 0.5, 1.0, 12.3, 6.0, arch, font_size=13,
    color=RGBColor(0x7D, 0xFF, 0xD3))

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Why Microservices
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_BG)
rect(s, 0, 0, 13.33, 1.0, DARK)
box(s, 0.5, 0.18, 12.3, 0.65, "Why Microservices Architecture?", font_size=30,
    bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Comparison table
headers = ["Criterion", "Monolith", "Microservices (Ours)"]
rows = [
    ["Scaling",       "Scale entire app",    "Scale only listing-service during peak"],
    ["Fault Isolation","One crash = all down","Auth crash → marketplace still works"],
    ["Deployment",    "Deploy everything",   "Deploy only changed service"],
    ["Team Work",     "Merge conflicts",     "Each service owned independently"],
    ["Technology",    "One stack only",      "Best tool per service"],
]
col_widths = [2.8, 3.5, 5.5]
col_x = [0.4, 3.3, 6.9]
for ci, (hdr, w, x) in enumerate(zip(headers, col_widths, col_x)):
    c = DARK if ci == 0 else (RED if ci == 1 else GREEN)
    rect(s, x, 1.2, w, 0.5, c)
    box(s, x + 0.05, 1.2, w - 0.1, 0.5, hdr, font_size=14, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    y = 1.75 + ri * 0.95
    for ci, (cell, w, x) in enumerate(zip(row, col_widths, col_x)):
        bg_c = RGBColor(0xF0, 0xF0, 0xF8) if ri % 2 == 0 else WHITE
        rect(s, x, y, w, 0.88, bg_c)
        tc = DARK if ci == 0 else (RED if ci == 1 else RGBColor(0x05, 0x7A, 0x55))
        box(s, x + 0.08, y + 0.1, w - 0.15, 0.7, cell, font_size=13,
            color=tc, bold=(ci == 0))

box(s, 0.5, 6.7, 12.3, 0.5,
    "Proven by Chaos Monkey test: stopping backend-auth-service-1 → marketplace, chat & monitoring keep running",
    font_size=13, color=ORANGE, align=PP_ALIGN.CENTER, bold=True)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Event-Driven Architecture
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
rect(s, 0, 0, 13.33, 1.0, RGBColor(0x0A, 0x0A, 0x1E))
box(s, 0.5, 0.18, 12.3, 0.65, "Event-Driven Architecture (Redis Pub/Sub)",
    font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

channels = [
    ("listing.event",       "listing-service", ["ai-service", "search-service"],
     "listing.created / listing.sold / listing.expired", ORANGE),
    ("audit.channel",       "ai-service",       ["admin-service"],
     "low_price_flag / high_price_flag / spam_rate_flag", RED),
    ("notification.channel","auth + listing +\nadmin services", ["notification-service"],
     "welcome.email / offer.accepted / listing.expired", GREEN),
    ("admin.event",         "admin-service",    ["listing-service"],
     "listing.removed / user.suspended", BLUE),
]
for i, (channel, pub, subs, events, col) in enumerate(channels):
    y = 1.2 + i * 1.45
    rect(s, 0.3, y, 12.7, 1.3, RGBColor(0x1A, 0x1A, 0x3E))
    rect(s, 0.3, y, 0.12, 1.3, col)
    box(s, 0.55, y + 0.05, 2.8, 0.4, channel, font_size=14, bold=True, color=col)
    box(s, 0.55, y + 0.45, 2.8, 0.7, "Publisher:\n" + pub, font_size=11, color=GRAY)
    box(s, 3.5,  y + 0.05, 4.5, 0.4, "Events:", font_size=11, color=GRAY)
    box(s, 3.5,  y + 0.4,  4.5, 0.75, events, font_size=12, color=WHITE)
    box(s, 8.2,  y + 0.05, 4.5, 0.4, "Subscribers:", font_size=11, color=GRAY)
    box(s, 8.2,  y + 0.4,  4.5, 0.75, " · ".join(subs), font_size=12, color=WHITE)

box(s, 0.5, 7.0, 12.3, 0.35,
    "Asynchronous communication — publishers don't wait for subscribers (loose coupling)",
    font_size=13, color=ORANGE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — AI Fraud Detection
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_BG)
rect(s, 0, 0, 13.33, 1.0, DARK)
box(s, 0.5, 0.18, 12.3, 0.65, "AI Fraud Detection Engine", font_size=30,
    bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 3 rules
for i, (icon, title, rule, example, col) in enumerate([
    ("📉", "LOW PRICE FLAG",  "Price < 60% of category reference",
     "Liquid Soap ref: 2,000 FCFA → flag if < 1,200 FCFA", RED),
    ("📈", "HIGH PRICE FLAG", "Price > 8× category reference",
     "Liquid Soap ref: 2,000 FCFA → flag if > 16,000 FCFA", ORANGE),
    ("🚨", "SPAM RATE FLAG",  "> 10 listings posted in 60 minutes",
     "Seller posts 12 listings in 1 hour → flagged", BLUE),
]):
    x = 0.35 + i * 4.3
    rect(s, x, 1.2, 4.1, 2.8, WHITE)
    rect(s, x, 1.2, 4.1, 0.5, col)
    box(s, x + 0.1, 1.2, 3.9, 0.45, icon + "  " + title, font_size=14,
        bold=True, color=WHITE)
    box(s, x + 0.1, 1.8, 3.9, 0.55, rule, font_size=14, bold=True, color=col)
    box(s, x + 0.1, 2.4, 3.9, 0.8, example, font_size=12, color=GRAY)

# Algorithm flow
rect(s, 0.3, 4.25, 12.7, 2.8, RGBColor(0x1A, 0x1A, 0x2E))
box(s, 0.5, 4.3, 12.3, 0.4,
    "How it works:", font_size=15, bold=True, color=ORANGE)
steps = [
    "1. listing-service publishes LISTING_CREATED event to Redis when new listing posted",
    "2. ai-service subscribes to listing.event channel — auto-triggers on every new listing",
    "3. ai-service gets category average price from Redis cache (6-hour TTL)",
    "4. Applies 3 rules → if flagged, publishes to audit.channel",
    "5. admin-service subscribes, saves to fraud_flags table in PostgreSQL",
    "6. ListingDetail page calls GET /api/admin/listing-flags/:id → shows red warning banner",
]
for i, step in enumerate(steps):
    box(s, 0.5, 4.75 + i * 0.35, 12.3, 0.33, step, font_size=12, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CI/CD Pipeline
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
rect(s, 0, 0, 13.33, 1.0, RGBColor(0x0A, 0x0A, 0x1E))
box(s, 0.5, 0.18, 12.3, 0.65, "Jenkins CI/CD Pipeline — 13 Stages", font_size=30,
    bold=True, color=WHITE, align=PP_ALIGN.CENTER)

stages = [
    ("1", "Checkout",       "Clone main branch from GitHub"),
    ("2", "Install Deps",   "npm ci for frontend + 6 services"),
    ("3", "Lint",           "ESLint — code style check"),
    ("4", "Run Tests",      "Jest — 232 tests across all services"),
    ("5", "Coverage",       "Generate lcov.info per service"),
    ("6", "SonarQube",      "Static analysis + security scan"),
    ("7", "Quality Gate",   "Fail pipeline if coverage < 85%"),
    ("8", "Build Images",   "docker build all 9 services"),
    ("9", "Push Images",    "Push to Docker Hub: praisesn/*"),
    ("10", "Deploy K8s",    "kubectl apply all manifests"),
    ("11", "Smoke Test",    "curl /health on all services"),
    ("12", "Notify",        "Send build result notification"),
    ("13", "Done",          "~10 minutes total"),
]
for i, (num, name, desc) in enumerate(stages):
    x = 0.3 + (i % 5) * 2.55
    y = 1.2 + (i // 5) * 2.05
    col = GREEN if num == "13" else (RED if num == "7" else ORANGE)
    rect(s, x, y, 2.3, 1.8, RGBColor(0x1A, 0x1A, 0x3E))
    rect(s, x, y, 2.3, 0.4, col)
    box(s, x + 0.05, y + 0.0, 2.2, 0.38, f"Stage {num}", font_size=12,
        bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, x + 0.05, y + 0.45, 2.2, 0.4, name, font_size=13, bold=True, color=col)
    box(s, x + 0.05, y + 0.88, 2.2, 0.78, desc, font_size=11, color=WHITE)

box(s, 0.5, 7.1, 12.3, 0.3,
    "Auto-triggers on every git push to main branch via GitHub Webhook → http://209.38.199.108:8090",
    font_size=12, color=ORANGE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Test Coverage
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_BG)
rect(s, 0, 0, 13.33, 1.0, DARK)
box(s, 0.5, 0.18, 12.3, 0.65, "Test Coverage — SonarQube Quality Gate",
    font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

services_cov = [
    ("auth-service",         "95%+", 42),
    ("user-service",         "95%+", 38),
    ("listing-service",      "95%+", 44),
    ("chat-service",         "95%+", 40),
    ("admin-service",        "95%+", 36),
    ("ai-service",           "95%+", 32),
]
for i, (svc, cov, tests) in enumerate(services_cov):
    x = 0.35 + (i % 3) * 4.2
    y = 1.2 + (i // 3) * 2.0
    rect(s, x, y, 3.9, 1.7, WHITE)
    rect(s, x, y, 3.9, 0.45, GREEN)
    box(s, x + 0.1, y + 0.0, 3.7, 0.42, svc, font_size=14, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    box(s, x + 0.1, y + 0.5, 3.7, 0.55, cov, font_size=28, bold=True,
        color=GREEN, align=PP_ALIGN.CENTER)
    box(s, x + 0.1, y + 1.05, 3.7, 0.45, f"{tests} tests", font_size=14,
        color=GRAY, align=PP_ALIGN.CENTER)

rect(s, 3.5, 5.5, 6.3, 1.6, DARK)
box(s, 3.6, 5.55, 6.1, 0.45, "TOTAL: 232 TESTS — 0 FAILING", font_size=18,
    bold=True, color=GREEN, align=PP_ALIGN.CENTER)
box(s, 3.6, 6.0,  6.1, 0.45, "SonarQube Quality Gate: PASSED", font_size=15,
    color=ORANGE, align=PP_ALIGN.CENTER)
box(s, 3.6, 6.45, 6.1, 0.4,
    "Framework: Jest + Supertest + Redis/PG Mocks", font_size=12,
    color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Docker & Kubernetes
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
rect(s, 0, 0, 13.33, 1.0, RGBColor(0x0A, 0x0A, 0x1E))
box(s, 0.5, 0.18, 12.3, 0.65, "Docker & Kubernetes Orchestration", font_size=30,
    bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Left: Docker
rect(s, 0.3, 1.2, 6.0, 5.9, RGBColor(0x1A, 0x1A, 0x3E))
box(s, 0.4, 1.25, 5.8, 0.45, "Docker Compose — 15 Containers",
    font_size=16, bold=True, color=ORANGE)
containers = [
    "backend-api-gateway-1     :8080",
    "backend-auth-service-1    :3001",
    "backend-listing-service-1 :3003",
    "backend-chat-service-1    :3004",
    "backend-ai-service-1      :3006",
    "backend-admin-service-1   :3005",
    "backend-postgres-1        :5432",
    "backend-redis-1           :6379",
    "backend-prometheus-1      :9090",
    "backend-grafana-1         :3009",
    "campustrade-jenkins        :8090",
    "campustrade-sonarqube      :9000",
]
for i, c in enumerate(containers):
    box(s, 0.4, 1.75 + i * 0.4, 5.8, 0.38, c, font_size=11,
        color=RGBColor(0x7D, 0xFF, 0xD3))

# Right: K8s
rect(s, 6.7, 1.2, 6.3, 5.9, RGBColor(0x1A, 0x1A, 0x3E))
box(s, 6.8, 1.25, 6.1, 0.45, "Kubernetes (k3s) + HPA",
    font_size=16, bold=True, color=BLUE)
hpa_items = [
    ("api-gateway-hpa",     "2 → 10 pods", "CPU > 60%"),
    ("auth-service-hpa",    "2 → 6 pods",  "CPU > 70%"),
    ("listing-service-hpa", "2 → 8 pods",  "CPU > 60%"),
    ("search-service-hpa",  "2 → 6 pods",  "CPU > 60%"),
    ("user-service-hpa",    "2 → 5 pods",  "CPU > 70%"),
]
for i, (name, scale, trigger) in enumerate(hpa_items):
    y = 1.75 + i * 1.0
    rect(s, 6.8, y, 6.0, 0.85, RGBColor(0x2A, 0x2A, 0x4E))
    box(s, 6.9, y + 0.05, 5.8, 0.35, name, font_size=13, bold=True, color=BLUE)
    box(s, 6.9, y + 0.42, 2.8, 0.35, scale,   font_size=12, color=GREEN)
    box(s, 9.9, y + 0.42, 2.8, 0.35, trigger, font_size=12, color=ORANGE)

box(s, 6.8, 6.85, 6.1, 0.3,
    "Self-healing: crashed pods restart automatically",
    font_size=12, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Monitoring
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_BG)
rect(s, 0, 0, 13.33, 1.0, DARK)
box(s, 0.5, 0.18, 12.3, 0.65, "Prometheus + Grafana Monitoring", font_size=30,
    bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Prometheus
rect(s, 0.3, 1.2, 5.9, 5.8, WHITE)
box(s, 0.4, 1.25, 5.7, 0.45, "Prometheus — http://209.38.199.108:9090",
    font_size=14, bold=True, color=RED)
box(s, 0.4, 1.75, 5.7, 0.35, "Scrapes every 15 seconds from all 9 services", font_size=12, color=GRAY)
metrics = [
    "http_requests_total        — requests by route/method/status",
    "http_request_duration_secs — response time histogram",
    "nodejs_heap_size_used      — JavaScript heap memory",
    "process_cpu_seconds_total  — CPU time used",
    "nodejs_eventloop_lag       — event loop latency",
    "nodejs_active_handles      — open connections",
]
for i, m in enumerate(metrics):
    box(s, 0.4, 2.2 + i * 0.7, 5.7, 0.6, m, font_size=12, color=DARK)

# Grafana
rect(s, 6.7, 1.2, 6.3, 5.8, WHITE)
box(s, 6.8, 1.25, 6.1, 0.45, "Grafana — http://209.38.199.108:3009",
    font_size=14, bold=True, color=ORANGE)
box(s, 6.8, 1.75, 6.1, 0.35, "admin / campustrade123", font_size=12, color=GRAY)
panels = [
    "HTTP Request Rate per Service (req/min)",
    "95th Percentile Response Time",
    "Node.js Heap Memory per Service",
    "CPU Usage per Service",
    "Total Active Connections",
    "Service Health Status (up/down)",
]
for i, p in enumerate(panels):
    rect(s, 6.8, 2.2 + i * 0.78, 6.0, 0.65, LIGHT_BG)
    box(s, 6.9, 2.25 + i * 0.78, 5.8, 0.55, "▪ " + p, font_size=13, color=DARK)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Results & Achievements
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
rect(s, 0, 0, 13.33, 1.0, RGBColor(0x0A, 0x1A, 0x0A))
box(s, 0.5, 0.18, 12.3, 0.65, "Results & Achievements", font_size=32,
    bold=True, color=GREEN, align=PP_ALIGN.CENTER)

achievements = [
    ("✅", "Live Application",     "http://209.38.199.108:4000 — 9 microservices running on DigitalOcean VPS"),
    ("✅", "232 Tests Passing",    "95%+ line coverage on all 6 tested services — SonarQube Quality Gate PASSED"),
    ("✅", "Full CI/CD Pipeline",  "13-stage Jenkins pipeline auto-triggers on every GitHub push"),
    ("✅", "Kubernetes Cluster",   "k3s with HPA autoscaling — 5 services configured to scale 2→10 pods"),
    ("✅", "AI Fraud Detection",   "3 rules: low price (<60%), high price (>800%), spam (>10 listings/60min)"),
    ("✅", "Prometheus + Grafana", "Real-time metrics from all 9 services — http://209.38.199.108:3009"),
    ("✅", "Chaos Monkey Ready",   "Stop any container — the other 14 keep running independently"),
    ("✅", "Ansible IaC",          "2 playbooks automate full server provisioning and deployment"),
]
for i, (icon, title, desc) in enumerate(achievements):
    y = 1.2 + i * 0.75
    rect(s, 0.3, y, 12.7, 0.68, RGBColor(0x0A, 0x2A, 0x0A))
    box(s, 0.4, y + 0.08, 0.5, 0.5, icon, font_size=16)
    box(s, 0.9, y + 0.04, 3.5, 0.35, title, font_size=14, bold=True, color=GREEN)
    box(s, 4.5, y + 0.04, 8.4, 0.55, desc, font_size=13, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Challenges & Solutions
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_BG)
rect(s, 0, 0, 13.33, 1.0, DARK)
box(s, 0.5, 0.18, 12.3, 0.65, "Challenges & Solutions", font_size=30,
    bold=True, color=WHITE, align=PP_ALIGN.CENTER)

challenges = [
    ("Jest Mock Leakage",       "Test values leaked between tests causing cascading failures",
     "Switched to resetAllMocks() + explicit mockImplementation restore in afterEach"),
    ("k3s Port Conflict",       "k3s Traefik ingress bound port 80, conflicting with Nginx",
     "Moved React frontend to port 4000 — free from Kubernetes service ports"),
    ("Coverage 54% Notify",     "notification-service event handlers never called in tests",
     "Implemented capturedHandler pattern: mock Redis subscribe, call callback directly"),
    ("Grafana Crash Loop",      "Grafana v12 broke Prometheus plugin provisioning",
     "Pinned to Grafana v11.4.0 and removed conflicting environment variables"),
    ("Docker Image Names",      "K8s manifests had campustrade/* but Docker Hub had praisesn/*",
     "Corrected all 9 Kubernetes YAML deployment files to use correct username"),
]
for i, (prob, cause, fix) in enumerate(challenges):
    y = 1.2 + i * 1.2
    rect(s, 0.3, y, 12.7, 1.1, WHITE)
    rect(s, 0.3, y, 0.15, 1.1, RED)
    box(s, 0.55, y + 0.05, 3.5, 0.38, "❌  " + prob, font_size=14, bold=True, color=RED)
    box(s, 4.2,  y + 0.05, 4.0, 0.35, "Cause: " + cause, font_size=11, color=GRAY)
    box(s, 0.55, y + 0.55, 3.5, 0.42, "✅  Fix:", font_size=13, bold=True, color=GREEN)
    box(s, 4.2,  y + 0.45, 8.6, 0.55, fix, font_size=12, color=DARK)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Recommendations
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
rect(s, 0, 0, 13.33, 1.0, RGBColor(0x0A, 0x0A, 0x1E))
box(s, 0.5, 0.18, 12.3, 0.65, "Recommendations — Future Improvements",
    font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

recs = [
    ("🔒", "HTTPS / SSL",              "Let's Encrypt via cert-manager — required for push notifications on all browsers"),
    ("💳", "Mobile Money Integration", "MTN Mobile Money + Orange Money (dominant in Cameroon) — in-app escrow payments"),
    ("🤖", "Machine Learning Pricing", "Replace rule-based with Random Forest/Gradient Boosting trained on transaction history"),
    ("📧", "Email Verification",       "Re-enable mandatory email verification to prevent spam accounts (disabled for demo)"),
    ("☁️",  "CDN for Images",           "Move image storage from VPS to Cloudflare R2 or AWS S3 — reduce costs, improve speed"),
    ("🏫", "Multi-Campus Expansion",   "Parameterise campus zones to support University of Buea, UYI and other campuses"),
]
for i, (icon, title, desc) in enumerate(recs):
    x = 0.35 + (i % 2) * 6.4
    y = 1.2 + (i // 2) * 2.0
    rect(s, x, y, 6.1, 1.75, RGBColor(0x1A, 0x1A, 0x3E))
    box(s, x + 0.1, y + 0.1, 0.7, 0.65, icon, font_size=26)
    box(s, x + 0.85, y + 0.1, 5.1, 0.45, title, font_size=15, bold=True, color=ORANGE)
    box(s, x + 0.85, y + 0.6, 5.1, 1.0,  desc,  font_size=12, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Conclusion
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
rect(s, 0, 0, 13.33, 1.0, RGBColor(0x0A, 0x0A, 0x1E))
box(s, 0.5, 0.18, 12.3, 0.65, "Conclusion", font_size=36, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)

box(s, 0.8, 1.15, 11.7, 0.65,
    "CampusTrade demonstrates a complete, production-grade software architecture",
    font_size=20, color=ORANGE, align=PP_ALIGN.CENTER, bold=True)

points = [
    "Event-Driven Microservices — 9 independent services communicating via Redis Pub/Sub",
    "Full DevOps Pipeline — Jenkins 13 stages: test → quality gate → build → deploy → monitor",
    "95.89% Test Coverage — 232 automated tests verified by SonarQube quality gate",
    "Cloud-Native — Docker + Kubernetes (k3s) with HPA autoscaling on DigitalOcean VPS",
    "AI Safety — 3-rule fraud detection engine with real-time buyer warnings",
    "Observable — Prometheus metrics + Grafana dashboards on all 9 services",
    "Resilient — Chaos Monkey proven: stop any service, the rest keep running",
]
for i, p in enumerate(points):
    box(s, 1.0, 1.95 + i * 0.65, 11.3, 0.6, "✓  " + p, font_size=15, color=WHITE)

rect(s, 1.5, 6.7, 10.3, 0.55, ORANGE)
box(s, 1.5, 6.7, 10.3, 0.55,
    "Live at: http://209.38.199.108:4000  |  GitHub: Nkinyampraises/smart-campus-market",
    font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Q&A
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
rect(s, 0, 0, 13.33, 0.12, ORANGE)
rect(s, 0, 7.38, 13.33, 0.12, ORANGE)

box(s, 1.0, 1.8, 11.3, 1.2, "Questions & Answers",
    font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(s, 1.0, 3.0, 11.3, 0.6,
    "Thank you for your attention",
    font_size=22, color=ORANGE, align=PP_ALIGN.CENTER)

for i, (label, url) in enumerate([
    ("Live App", "http://209.38.199.108:4000"),
    ("Jenkins",  "http://209.38.199.108:8090"),
    ("SonarQube","http://209.38.199.108:9000"),
    ("Grafana",  "http://209.38.199.108:3009"),
]):
    x = 0.9 + i * 2.9
    rect(s, x, 3.85, 2.6, 1.0, RGBColor(0x1A, 0x1A, 0x3E))
    box(s, x + 0.05, 3.9, 2.5, 0.4, label, font_size=13, bold=True,
        color=ORANGE, align=PP_ALIGN.CENTER)
    box(s, x + 0.05, 4.3, 2.5, 0.45, url.replace("http://", ""),
        font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

box(s, 1.0, 5.3, 11.3, 0.4,
    "Praises Ncha — nkinyampraises.ncha@ictuniversity.edu.cm",
    font_size=14, color=GRAY, align=PP_ALIGN.CENTER)
box(s, 1.0, 5.75, 11.3, 0.4,
    "Kongyu Jesse Ntani — ICTU20234195",
    font_size=14, color=GRAY, align=PP_ALIGN.CENTER)
box(s, 1.0, 6.25, 11.3, 0.4,
    "Instructor: Engr. TEKOH PALMA  |  SEN3244 Software Architecture  |  ICT University 2026",
    font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────────────────────
out = r"C:\Users\MATRIX TECHNOLOGY\Documents\smart campus market\CampusTrade_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
